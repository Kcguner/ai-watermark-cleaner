# tests/test_edge_cases.py
"""Edge-case tests (group D). New file only; existing tests untouched."""
import io

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app import create_app
from cleaners.image_cleaner import clean_image
from cleaners.pptx_cleaner import clean_pptx
from cleaners.registry import get_signature


def _client():
    return create_app().test_client()


def _png_bytes(size=(100, 100), color="white"):
    buf = io.BytesIO()
    Image.new("RGB", size, color).save(buf, format="PNG")
    buf.seek(0)
    return buf


def test_corrupt_png_returns_400():
    c = _client()
    data = {"file": (io.BytesIO(b"this is not a png at all"), "bad.png"), "source": "gamma"}
    r = c.post("/api/clean", data=data, content_type="multipart/form-data")
    assert r.status_code == 400


def test_unsupported_txt_returns_400():
    c = _client()
    data = {"file": (io.BytesIO(b"hello world"), "note.txt"), "source": "gamma"}
    r = c.post("/api/clean", data=data, content_type="multipart/form-data")
    assert r.status_code == 400


def test_missing_file_returns_400():
    c = _client()
    r = c.post("/api/clean", data={"source": "gamma"}, content_type="multipart/form-data")
    assert r.status_code == 400


def test_unknown_source_falls_back_generic_200():
    c = _client()
    data = {"file": (_png_bytes(), "t.png"), "source": "definitely-not-a-source-xyz"}
    r = c.post("/api/clean", data=data, content_type="multipart/form-data")
    assert r.status_code == 200
    assert r.content_type.startswith("image/png")


def test_clean_white_image_gamma_unchanged_size():
    img = Image.new("RGB", (200, 200), "white")
    sig = get_signature("gamma")
    out = clean_image(img, sig)
    assert out.size == img.size


def test_pptx_top_title_kept_bottom_watermark_removed(tmp_path):
    src = tmp_path / "in.pptx"
    out = tmp_path / "out.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    top_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(5), Inches(1))
    top_box.text_frame.text = "Gamma project"
    bot_top = int(prs.slide_height * 0.85)
    bot_box = slide.shapes.add_textbox(Inches(0.5), bot_top, Inches(5), Inches(0.5))
    bot_box.text_frame.text = "Made with Gamma watermark"
    prs.save(str(src))

    report = clean_pptx(str(src), str(out), get_signature("gamma"))
    assert out.exists()
    assert report["removed"] >= 1
    cleaned = Presentation(str(out))
    texts = [s.text for sl in cleaned.slides for s in sl.shapes if s.has_text_frame]
    assert any("Gamma project" in t for t in texts)
    assert not any("Made with Gamma" in t for t in texts)
