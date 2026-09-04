"""Test fixture generator for Gamma watermark samples.

Creates synthetic English fixtures with a bottom "Made with Gamma" band:
- sample_gamma.png (800x600)
- sample_gamma.pdf (595x842, A4-ish)
- sample_gamma.pptx (blank slide + watermark textbox)
"""
from pathlib import Path

FIXTURES_DIR = Path(__file__).resolve().parent

IMAGE_NAME = "sample_gamma.png"
PDF_NAME = "sample_gamma.pdf"
PPTX_NAME = "sample_gamma.pptx"

SAMPLE_TEXT = "Sample content"
WATERMARK_TEXT = "Made with Gamma"


def make_image(path: Path | str = FIXTURES_DIR / IMAGE_NAME) -> Path:
    """Create 800x600 white PNG with content and bottom gray watermark band."""
    from PIL import Image, ImageDraw

    path = Path(path)
    width, height = 800, 600
    img = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(img)

    # Black rect outline to simulate content area.
    draw.rectangle([50, 50, width - 50, height - 120], outline="black", width=2)
    # Main content text.
    draw.text((70, 80), SAMPLE_TEXT, fill="black")
    # Bottom 7% gray band.
    band_height = int(height * 0.07)
    band_top = height - band_height
    draw.rectangle([0, band_top, width, height], fill=(245, 245, 245))
    # Watermark text inside band.
    draw.text((20, band_top + 10), WATERMARK_TEXT, fill=(130, 130, 130))

    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path)
    return path


def make_pdf(path: Path | str = FIXTURES_DIR / PDF_NAME) -> Path:
    """Create A4-ish 595x842 PDF with content and bottom watermark band."""
    import fitz

    path = Path(path)
    width, height = 595, 842
    doc = fitz.open()
    page = doc.new_page(width=width, height=height)

    # Main content text.
    page.insert_text((72, 100), SAMPLE_TEXT, fontsize=18, color=(0, 0, 0))

    # Bottom 7% band rect fill (245/255 gray).
    band_top = height * 0.93
    band_rect = fitz.Rect(0, band_top, width, height)
    gray = 245 / 255
    page.draw_rect(band_rect, color=None, fill=(gray, gray, gray))

    # Watermark text, 9pt gray.
    page.insert_text(
        (72, band_top + 25), WATERMARK_TEXT, fontsize=9, color=(0.5, 0.5, 0.5)
    )

    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(path)
    doc.close()
    return path


def make_pptx(path: Path | str = FIXTURES_DIR / PPTX_NAME) -> Path:
    """Create PPTX with content textbox and bottom watermark textbox."""
    from pptx import Presentation
    from pptx.util import Inches

    path = Path(path)
    prs = Presentation()
    # Blank layout: index 6 in default template, fall back to 5.
    layouts = prs.slide_layouts
    blank = layouts[6] if len(layouts) > 6 else layouts[5]
    slide = prs.slides.add_slide(blank)

    # Main content textbox.
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(2)
    )
    content_box.text_frame.text = SAMPLE_TEXT

    # Bottom watermark textbox named "Gamma Watermark".
    wm_box = slide.shapes.add_textbox(
        Inches(0.5),
        prs.slide_height - Inches(1.0),
        prs.slide_width - Inches(1.0),
        Inches(0.6),
    )
    wm_box.text_frame.text = WATERMARK_TEXT
    wm_box.name = "Gamma Watermark"

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


def main() -> None:
    make_image()
    make_pdf()
    make_pptx()
    print("fixtures written")


if __name__ == "__main__":
    main()
