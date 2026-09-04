from pathlib import Path
from pptx import Presentation
from cleaners.pptx_cleaner import clean_pptx
from cleaners.registry import get_signature

FIX = Path(__file__).parent / "fixtures" / "sample_gamma.pptx"

def test_gamma_shape_removed(tmp_path):
    out = tmp_path / "clean.pptx"
    report = clean_pptx(str(FIX), str(out), get_signature("gamma"))
    assert out.exists()
    assert report["removed"] >= 1
    prs = Presentation(out)
    texts = [s.text for sl in prs.slides for s in sl.shapes if s.has_text_frame]
    assert not any("Gamma" in t for t in texts)
