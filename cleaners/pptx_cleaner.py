"""PPTX cleaner: removes watermark-like shapes."""
from __future__ import annotations
from pptx import Presentation

KEYWORDS = ("gamma", "gemini", "watermark", "nano banana", "made with")


def _iter_shapes(shapes):
    """Yield shapes recursively, descending into group shapes."""
    for s in shapes:
        yield s
        try:
            is_group = s.shape_type == 6
        except Exception:
            is_group = False
        if is_group and hasattr(s, "shapes"):
            try:
                yield from _iter_shapes(s.shapes)
            except Exception:
                continue


def _is_bottom(shape, slide_height) -> bool:
    """Return True iff the shape sits in the bottom 25% of the slide."""
    try:
        if slide_height is None or not slide_height:
            return False
        top = shape.top
        if top is None:
            return False
        return (top / slide_height) > 0.75
    except Exception:
        return False


def _shape_text(shape) -> str:
    """Best-effort text extraction for text frames and tables."""
    texts = []
    if getattr(shape, "has_text_frame", False):
        try:
            texts.append(shape.text or "")
        except Exception:
            pass
    if getattr(shape, "has_table", False):
        try:
            for row in shape.table.rows:
                for cell in row.cells:
                    try:
                        texts.append(cell.text or "")
                    except Exception:
                        continue
        except Exception:
            pass
    return "\n".join(texts)


def _is_watermark(shape, slide_height) -> bool:
    name = (getattr(shape, "name", "") or "").lower()
    if "watermark" in name or "gamma" in name or "gemini" in name:
        # Name match alone is not enough; require bottom position
        # to avoid deleting legit titled shapes.
        if _is_bottom(shape, slide_height):
            return True
        return False
    txt = _shape_text(shape).lower()
    if txt and any(k in txt for k in KEYWORDS):
        # Only remove if in the bottom area (protects main content headings).
        if _is_bottom(shape, slide_height):
            return True
    return False


def clean_pptx(in_path: str, out_path: str, signature: dict) -> dict:
    prs = Presentation(in_path)
    slide_height = prs.slide_height
    removed = 0
    for slide in prs.slides:
        for shape in list(_iter_shapes(slide.shapes)):
            if _is_watermark(shape, slide_height):
                sp = shape._element
                sp.getparent().remove(sp)
                removed += 1
    # Clear core properties
    prs.core_properties.author = ""
    prs.core_properties.last_modified_by = ""
    prs.save(out_path)
    return {"slides": len(prs.slides), "removed": removed}
